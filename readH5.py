import h5py
import numpy as np
import matplotlib.pyplot as plt
from matplotlib.transforms import Bbox # 用于获取axes边界
from matplotlib.ticker import MaxNLocator # 用于设置tick最大刻度数
import matplotlib.colors as mcolors
import re # 用于排序
from scipy.optimize import curve_fit  # 用于拟合
import pprint # 用于分行打印
import os # 用于判断文件是否存在
"""PPT"""
from pptx import Presentation
from pptx.util import Inches, Cm, Pt

# 5
nums = [5]
powers = [0.01, 0.02, 0.05, 0.1, 0.15, 0.2, 0.3, 0.4]  # mW
thumb_images = [0, 1]
# 6 10 11 12 13
# nums = [6, 10, 11, 12, 13]
# powers = [0.01, 0.02, 0.05, 0.1, 0.15, 0.2, 0.3, 0.4, 0.6, 0.8, 1]  # mW
# thumb_images = [0, 1]


"""定义高斯函数。用于拟合"""
def gaussian(x, height, center, width, shift):
    return height * np.exp(-(x - center) ** 2 / (2 * width ** 2)) + shift


"""获取axes的边界。用于保存subplot"""
def full_extent(ax, pad=0.0):
    """Get the full extent of an axes, including axes labels, tick labels, and titles."""
    # For text objects, we need to draw the figure first, otherwise the extents
    # are undefined.
    ax.figure.canvas.draw()
    # items = [ax]
    items = ax.get_xticklabels()
    items += ax.get_yticklabels()
    items += [ax.title]
    items += [ax.xaxis.label, ax.yaxis.label]
    bbox_items = [item.get_window_extent() for item in items if item.get_visible()]
    bbox_plot = ax.get_window_extent()
    bbox = Bbox.union(bbox_items + [bbox_plot])
    return bbox.expanded(1.0+pad, 1.0+pad)


"""save subplot"""
def save_subfig(fig, ax, save_name_full):
    bbox = ax.get_tightbbox(fig.canvas.get_renderer()).expanded(1.02, 1.02)
    extent = bbox.transformed(fig.dpi_scale_trans.inverted())
    fig.savefig(save_name_full, bbox_inches=extent)


"""打开h5文件，打印一级目录key"""
file_path = r'C:\Users\a1033\Desktop\临时备份\20240620\2024-05-31.h5'
f = h5py.File(file_path, "r")


"""取指定粒子们（under一级目录）"""
particles = []
for key in f['/NPoM_BPT_4'].keys():
    # if 'Particle' in key and 'ParticleScannerScan' not in key:
    if 'kinetic_SERS' in key:
        particles.append(key)


"""对粒子们的根据名称排序"""
particles.sort(key=lambda i: int(re.search(r'(\d+)', i).group()) if re.search(r'(\d+)', i) else float('inf'))
print('Preprocession==========completed')

print(particles)

"""处理DF_scan——始"""
# # 逐个打开particle中的DF_scan,并作图
# nFig = 0
# nSubplot = 0
# i = 0
# j = 0
# for key in particles:
#     try:
#         DF_scan = f[f'/{key}/DF_scan']
#         its = np.array(DF_scan)
#         # 重构DF光谱，每列取最大值，组成plot
#         wav = DF_scan.attrs['wavelengths']
#         bgd = DF_scan.attrs['background']
#         ref = DF_scan.attrs['reference']
#         np.seterr(divide='ignore', invalid='ignore')  # 消除被除数为0的警告
#         its_c = (its - bgd) / (ref - bgd)  # its_n:intensity combined
#         # 指定波段输出plot
#         max_its = np.max(its_c, axis=0)
#         idx_min = np.argmin(abs(wav-400))
#         idx_max = np.argmin(abs(wav-1000))
#         L = idx_max-idx_min
#         x = wav[idx_min:idx_max]
#         y = max_its[idx_min:idx_max]
#         # 作图：重构后的DF_scan plot
#         if (nSubplot % 12) == 0:
#             Fig, ax = plt.subplots(3, 4, figsize=(12, 8), tight_layout=True)
#             # 保存subplot的时候应用tight_layout严格防止子图重叠
#             nFig += 1
#         else:
#             pass
#         i = (nSubplot % 12) // 4
#         j = (nSubplot % 12) % 4
#         ax[i][j].plot(x, y, linewidth=0.5)
#         ax[i, j].tick_params(axis='both', which='major', labelsize=8)  # 设置主刻度字体大小
#         ax[i, j].tick_params(axis='both', which='minor', labelsize=8)  # 设置次刻度字体大小
#         # 使用MaxNLocator设置刻度数量
#         ax[i, j].xaxis.set_major_locator(MaxNLocator(nbins=6))  # 设置x轴刻度的最大数量为6
#         ax[i, j].yaxis.set_major_locator(MaxNLocator(nbins=6))  # 设置y轴刻度的最大数量为6
#         ax[i, j].set_xlabel('Wavelength(nm)', fontsize=8)
#         ax[i, j].set_ylabel('Normalized its', fontsize=8)
#         ax[i, j].set_title(f'{key}---DF_scan', fontsize=10)
#         extent = full_extent(ax[i, j]).transformed(Fig.dpi_scale_trans.inverted())
#         # savepath = r'C:\Users\a1033\Desktop\临时备份\20240620\Particles'
#         # savename = savepath + fr'\{key}-DF_scan.png'
#         # plt.savefig(savename, bbox_inches=extent)
#         # print(f"add file: {key}-DF_scan.png")
#         if key == 'Particle_5_bgd':
#             DF_Particle_5_bgd = y
#             wav_Particle_5 = x
#         elif key == 'Particle_5':
#             DF_Particle_5 = y
#         else:
#             print('error: DF_Particle_5(_bgd)')
#         nSubplot += 1
#     except KeyError as e:
#         print(f"Error: {e}")
# # bgd 5
# Fig = plt.figure()
# ax = Fig.add_subplot(111)
# ax.plot(wav_Particle_5, DF_Particle_5 - DF_Particle_5_bgd)
# ax.tick_params(axis='both', which='major', labelsize=8)  # 设置主刻度字体大小
# ax.tick_params(axis='both', which='minor', labelsize=8)  # 设置次刻度字体大小
# ax.xaxis.set_major_locator(MaxNLocator(nbins=6))  # 设置x轴刻度的最大数量为6
# ax.yaxis.set_major_locator(MaxNLocator(nbins=6))  # 设置y轴刻度的最大数量为6
# ax.set_xlabel('Wavelength(nm)', fontsize=8)
# ax.set_ylabel('Normalized intensity(counts)', fontsize=8)
# ax.set_title(f'DF_scan-Minus bgd', fontsize=10)
# savepath = r'C:\Users\a1033\Desktop\临时备份\20240620\Particles'
# savename = savepath + fr'\particle5_minus_bgd.png'
# plt.savefig(savename, bbox_inches='tight', dpi=300)
"""处理DF_scan——末"""

"""处理kinetic_SERS——始"""
SERS_bgd = f['/ParticleScannerScan_6/Particle_92/kinetic_SERS_0.1mW']
its_bgd = np.mean(SERS_bgd, axis=0)
"初始化"
nFig = 0
nSubplot = 0
i = 0
j = 0
stat = np.zeros(0)
nStat = 0
for key in particles:
    try:
        "读取数据"
        # SERS = f[f'/ParticleScannerScan_7/{key}/kinetic_SERS']
        SERS = f[f'/NPoM_BPT_4/{key}']
        its = np.array(SERS)
        wav = SERS.attrs['wavelengths']
        "设置每页最大子图数目"
        if (nSubplot % 12) == 0:
            Fig, ax = plt.subplots(3, 4, figsize=(15, 10), tight_layout=True)
            # 保存subplot的时候应用tight_layout严格防止子图重叠
            nFig += 1
        else:
            pass

        "plot SERS平均谱"
        x = wav
        y = np.mean(its, axis=0)
        y = y - its_bgd
        i = (nSubplot % 12) // 4
        j = (nSubplot % 12) % 4  # 子图横纵索引
        ax[i][j].plot(x, y, linewidth=0.5)
        ax[i, j].tick_params(axis='both', which='major', labelsize=8)  # 设置主刻度字体大小
        ax[i, j].tick_params(axis='both', which='minor', labelsize=8)  # 设置次刻度字体大小
        ax[i, j].xaxis.set_major_locator(MaxNLocator(nbins=6))  # 设置x轴刻度的最大数量为6
        ax[i, j].yaxis.set_major_locator(MaxNLocator(nbins=6))  # 设置y轴刻度的最大数量为6
        ax[i, j].set_xlabel('Wavelength(nm)', fontsize=8)
        ax[i, j].set_ylabel('Avarage SERS intensity', fontsize=8)
        ax[i, j].set_title(f'{key}---SERS', fontsize=10)
        nSubplot += 1

        """指定波段峰拟合(仅拟合一条)"""
        try:
            x = np.flipud(wav)
            y = np.flipud(y)
            idx_fitMin = np.argmin(abs(x - 695))
            idx_fitMax = np.argmin(abs(x - 710))  # 取定拟合区间
            popt, pcov = curve_fit(gaussian, x[idx_fitMin:idx_fitMax], y[idx_fitMin:idx_fitMax], p0=[10000, 703, 2, 0],
                                   method='trf')  # 获取拟合参数列表popt: height, center, width, shift
            ax[i][j].plot(x[idx_fitMin:idx_fitMax], gaussian(x[idx_fitMin:idx_fitMax], *popt), color='blue',
                          label=popt[0])
            ax[i][j].text(0.8, 0.85, "{:.1f}".format(popt[0] + popt[3]), transform=ax[i, j].transAxes, fontsize=15,
                          ha='center', color='red')
            # if popt[0] + popt[3] > 0:
            #     stat = np.append(stat, popt[0] + popt[3])  # 统计SERS峰值
            stat = np.append(stat, popt[0] + popt[3])  # 统计SERS峰值
            nStat += 1

            "save figure"
            # extent = full_extent(ax[i, j]).transformed(Fig.dpi_scale_trans.inverted())
            # savepath = r'C:\Users\a1033\Desktop\临时备份\20240620\Particles'
            # savename = savepath + fr'\{key}-SERS.png'
            # plt.savefig(savename, bbox_inches=extent)
            # print(f"add file: {key}-SERS.png")
        except Exception as e:
            print(f'fitting error this file: {key} >>>>> ', end='')
            print(e)
            pass
    except KeyError as e:
        print(f'key error this file: {key} >>>>> ', end='')
        print(e)

"""直方图"""
# Fig = plt.figure()
# ax = Fig.add_subplot(111)
# freqs, bins, patches = ax.hist(stat, bins=30, alpha=0.75, color='blue', edgecolor='black')
# ax.tick_params(axis='both', which='major', labelsize=8)  # 设置主刻度字体大小
# ax.tick_params(axis='both', which='minor', labelsize=8)  # 设置次刻度字体大小
# ax.xaxis.set_major_locator(MaxNLocator(nbins=6))  # 设置x轴刻度的最大数量为6
# ax.yaxis.set_major_locator(MaxNLocator(nbins=6))  # 设置y轴刻度的最大数量为6
# ax.set_xlabel('SERS intensity', fontsize=8)
# ax.set_ylabel('Frequency', fontsize=8)
# ax.set_title(f'Histogram of SERS', fontsize=10)
# # 添加数值标签
# for freq, rect in zip(freqs, patches):
#     height = rect.get_height()
#     # 计算条形的中心位置
#     center_x = rect.get_x() + rect.get_width() / 2
#     # 显示每个直方中心的横坐标位置和高度
#     if not freq == 0:
#         ax.text(center_x + 0.18 * rect.get_width(), height + 0.02 * max(freqs), f'{center_x:.1f}', color='purple', ha='center', va='bottom', fontsize=7,  rotation=90)
# savepath = r'C:\Users\a1033\Desktop\临时备份\20240620\NPoM_BPT_4'
# savename = savepath + fr'\Histogram_SERS.png'
# plt.savefig(savename, bbox_inches='tight', dpi=300)


"""x-I图"""
x = len(stat)
x = range(x)
Fig = plt.figure()
ax = Fig.add_subplot(111)
ax.plot(x, stat)
ax.tick_params(axis='both', which='major', labelsize=8)  # 设置主刻度字体大小
ax.tick_params(axis='both', which='minor', labelsize=8)  # 设置次刻度字体大小
ax.xaxis.set_major_locator(MaxNLocator(nbins=6))  # 设置x轴刻度的最大数量为6
ax.yaxis.set_major_locator(MaxNLocator(nbins=6))  # 设置y轴刻度的最大数量为6
ax.set_xlabel('x position', fontsize=8)
ax.set_ylabel('SERS Intensity', fontsize=8)
ax.set_title(f'x-I(SERS)', fontsize=10)
ax.legend(['NPoM_BPT_4'])
# 拟合
try:
    x = x
    y = stat
    popt, pcov = curve_fit(gaussian, x, y, p0=[500, 6, 2, 0],
                           method='trf')  # 获取拟合参数列表popt: height, center, width, shift
    xmin = x[0]
    xmax = x[:-1]
    x = np.linspace(xmin, xmax, 200)
    ax.plot(x, gaussian(x, *popt), color='pink',
                  label=popt[0])
    ax.text(0.15, 0.85, "width: {:.3f}".format(2 * np.sqrt(2 * np.log(2)) * popt[2]), transform=ax.transAxes, fontsize=15,
                  ha='center', color='red')
    savepath = r'C:\Users\a1033\Desktop\临时备份\20240620\NPoM_BPT_4'
    savename = savepath + fr'\x-I(SERS).png'
    plt.savefig(savename, bbox_inches='tight', dpi=300)
except Exception as e:
    print(e)
    pass

"""处理kinetic_SERS——末"""


"""processing——末"""

    # for image in thumb_images:
    #     """Image——始
    #     """
    #     nFig = 0
    #     nSubplot = 0
    #     i = 0
    #     j = 0
    #     for key in particles:
    #         try:
    #             NP_image = f[f'/ParticleScannerScan_{num}/{key}/CWL.thumb_image_{image}']
    #             NP_image = np.array(NP_image)
    #             if (nSubplot % 12) == 0:
    #                 Fig, ax = plt.subplots(3, 4, figsize=(12, 8), tight_layout=True)
    #                 # 保存subplot的时候应用tight_layout严格防止子图重叠
    #                 nFig += 1
    #             else:
    #                 pass
    #             i = (nSubplot % 12) // 4
    #             j = (nSubplot % 12) % 4
    #             ax[i][j].imshow(NP_image)
    #             ax[i, j].set_title(f'{key}-CWL.thumb_image_{image}')
    #             nSubplot += 1
    #             extent = full_extent(ax[i, j]).transformed(Fig.dpi_scale_trans.inverted())
    #             savename_image = savepath_image + f'/{key}-CWL.thumb_image_{image}.png'
    #             plt.savefig(savename_image, bbox_inches=extent)
    #         except KeyError as e:
    #             print(e)
    #             passpt KeyError as e:
    #     #             print(e)
    #     #             pass
    #     # plt.show()
    #     """Image——末
    #     """

print('All Processing==========completed')

    # """PPT——始
    # """
    # """单张"""
    # # # 获取包含数据的particles路径（含文件名）
    # # particles_valid = []
    # # for key in particles:
    # #     try:
    # #         SERS = f[f'/{key}/kinetic_SERS']
    # #         particles_valid.append(fr'C:\Users\a1033\Desktop\临时备份\20240620\Particles-SERS\SERS_all\{key}-SERS.png')
    # #     except Exception as e:
    # #         print(e)
    # #         pass
    # # pprint.pprint(particles_valid)
    # # # 创建一个新的演示文稿对象
    # # prs = Presentation()
    # # prs.slide_height = Cm(19.05) # 设置页面高度
    # # prs.slide_width = Cm(33.867) # 设置页面宽度
    # #
    # # # 设定每页图片排列的行数和列数
    # # rows = 3
    # # cols = 4
    # #
    # # # 图片路径列表
    # # image_paths = particles_valid
    # #
    # # # 每张幻灯片中的图片数量
    # # images_per_slide = rows * cols
    # #
    # # # 图片的大小和间距（单位：英寸）
    # # image_width = Cm(7.56)
    # # image_height = Cm(5.67)
    # # # h_spacing = Cm(0.1)
    # # # v_spacing = Cm(0.1)
    # # h_spacing = Cm(0)
    # # v_spacing = Cm(0)
    # #
    # # # 创建并添加图片到演示文稿
    # # for i, image_path in enumerate(image_paths):
    # #     # 每当i是images_per_slide的倍数时，添加一个新幻灯片
    # #     if i % images_per_slide == 0:
    # #         slide = prs.slides.add_slide(prs.slide_layouts[5])  # 使用空白布局
    # #
    # #     # 计算图片在当前幻灯片中的位置
    # #     row = (i % images_per_slide) // cols
    # #     col = (i % images_per_slide) % cols
    # #     left = Inches(0.5) + col * (image_width + h_spacing)
    # #     top = Inches(0.5) + row * (image_height + v_spacing)
    # #
    # #     # 添加图片到当前幻灯片
    # #     slide.shapes.add_picture(image_path, left, top, width=image_width, height=image_height)
    # #
    # # # 保存演示文稿
    # # prs.save(fr'C:\Users\a1033\Desktop\临时备份\20240620\Particles-SERS\SERS_all\presentation.pptx')
    #
    # """对比"""
    # # 创建一个新的演示文稿对象
    # prs = Presentation()
    # prs.slide_height = Cm(19.05) # 设置页面高度
    # prs.slide_width = Cm(33.867) # 设置页面宽度
    # # 设定每页图片排列的行数和列数
    # rows = 14
    # cols = 18
    #
    # # 每张幻灯片中的图片数量
    # images_per_slide = rows * cols
    #
    # # 图片的大小和间距（单位：英寸）
    # # image_width = Cm(7.56)
    # # image_height = Cm(5.67)
    # image_width = Cm(4)
    # image_height = Cm(3)
    # h_spacing = Cm(0)
    # v_spacing = Cm(0)
    # i = 0
    # rows_image = [0, 1]
    # rows_SERS = range(3, 14)
    # for key in particles:
    #     # 每当i是每张幻灯片列数的倍数时，添加一个新幻灯片
    #     if i % cols == 0:
    #         slide = prs.slides.add_slide(prs.slide_layouts[6])  # 使用空白布局
    #         items = []
    #     try:
    #         for image, row_image in zip(thumb_images, rows_image):
    #             image_path = savepath_image + f'/{key}-CWL.thumb_image_{image}.png'
    #             # 计算图片在当前幻灯片中的位置
    #             row = row_image
    #             print(row)
    #             col = i % cols
    #             left = Inches(0) + col * (image_width + h_spacing)
    #             top = Inches(0) + row * (image_height + v_spacing)
    #             # 添加图片到当前幻灯片
    #             slide.shapes.add_picture(image_path, left, top, width=image_width, height=image_height)
    #     except Exception as e:
    #         print(e)
    #         # print('NP_image_0 KeyError!')
    #         pass
    #
    #     try:
    #         DF_path = savepath_DF + f'/{key}-DF_zscan.png'
    #         # 计算图片在当前幻灯片中的位置
    #         row = 2
    #         print(row)
    #         col = i % cols
    #         left = Inches(0) + col * (image_width + h_spacing)
    #         top = Inches(0) + row * (image_height + v_spacing)
    #         # 添加图片到当前幻灯片
    #         slide.shapes.add_picture(DF_path, left, top, width=image_width, height=image_height)
    #     except Exception as e:
    #         print(e)
    #         # print('DF_scan KeyError!')
    #         pass
    #
    #     try:
    #         for power, row_SERS in zip(powers, rows_SERS):
    #             SERS_path = savepath_SERS + f'/{key}-SERS_{power}mW.png'
    #             # 计算图片在当前幻灯片中的位置
    #             row = row_SERS
    #             print(row)
    #             col = i % cols
    #             left = Inches(0) + col * (image_width + h_spacing)
    #             top = Inches(0) + row * (image_height + v_spacing)
    #             # 添加图片到当前幻灯片
    #             slide.shapes.add_picture(SERS_path, left, top, width=image_width, height=image_height)
    #             slide.shapes.add_group_shape()
    #     except Exception as e:
    #         print(e)
    #         # print('kinetic_SERS KeyError!')
    #         pass
    #     print("Add filenow to PPT==============completed!")
    #     i += 1
    # # 保存演示文稿
    # prs.save(fr'C:\Users\a1033\Desktop\临时备份\20240620\ParticleScannerScan_{num}\Particles.pptx')
    # """PPT——末
    # """


"show pictures"
plt.show()